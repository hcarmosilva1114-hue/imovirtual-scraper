#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Imovirtual → CSV + PowerPoint (Template Hugo Silva)
Autor: preparado para Hugo Silva (RE/MAX Oceanus)

AVISO LEGAL: Usa apenas dados/fotos com autorização. O scraping pode violar Termos de Uso.
"""

import asyncio
import argparse
import csv
import io
import json
import re
import time
from pathlib import Path
from typing import Dict, Any, List, Optional

import pandas as pd
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_FILE = "template_hugosilva.pptx"

# ---------------------------
# Navegação com Playwright
# ---------------------------
async def fetch_html(url: str, render: str = "networkidle", timeout_ms: int = 45000) -> str:
    from playwright.async_api import async_playwright
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(user_agent=(
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36"
        ))
        page = await context.new_page()
        await page.goto(url, wait_until=render, timeout=timeout_ms)
        html = await page.content()
        await browser.close()
        return html

# ---------------------------
# Parsing auxiliar
# ---------------------------
def text_or_none(el) -> Optional[str]:
    return el.get_text(strip=True) if el else None

def find_label_value(soup: BeautifulSoup, labels: List[str]) -> Optional[str]:
    label_pattern = re.compile(r"^\s*(%s)\s*:?$" % "|".join([re.escape(x) for x in labels]), re.I)
    for dt in soup.find_all("dt"):
        if dt.string and label_pattern.match(dt.get_text(strip=True)):
            dd = dt.find_next("dd")
            if dd:
                return dd.get_text(" ", strip=True)
    for li in soup.find_all("li"):
        strong = li.find(["strong", "span"])
        if strong and label_pattern.match(strong.get_text(strip=True)):
            txt = li.get_text(" ", strip=True)
            lab = strong.get_text(strip=True)
            val = txt.replace(lab, "").strip(" :\u00a0-")
            if val:
                return val
    for node in soup.find_all(["div", "span"]):
        if node.get_text(strip=True) and label_pattern.match(node.get_text(strip=True)):
            sib = node.find_next_sibling()
            if sib and sib.get_text(strip=True):
                return sib.get_text(" ", strip=True)
    return None

def parse_json_ld(soup: BeautifulSoup) -> Dict[str, Any]:
    data: Dict[str, Any] = {}
    for tag in soup.find_all("script", type="application/ld+json"):
        try:
            obj = json.loads(tag.string or "{}")
        except Exception:
            continue
        if isinstance(obj, list):
            for x in obj:
                _merge_realestate(data, x)
        else:
            _merge_realestate(data, obj)
    return data

def _merge_realestate(dst: Dict[str, Any], obj: Dict[str, Any]):
    if not isinstance(obj, dict):
        return
    if "@graph" in obj and isinstance(obj["@graph"], list):
        for g in obj["@graph"]:
            _merge_realestate(dst, g)
    offer = obj.get("offers") or {}
    address = obj.get("address") or {}
    dst.setdefault("title", obj.get("name"))
    dst.setdefault("description", obj.get("description"))
    if isinstance(offer, dict):
        price = offer.get("price") or offer.get("lowPrice")
        dst.setdefault("price", price)
        currency = offer.get("priceCurrency")
        if price and currency:
            dst["price_str"] = f"{price} {currency}"
    if isinstance(address, dict):
        loc = " ".join([
            str(address.get("addressLocality") or ""),
            str(address.get("addressRegion") or ""),
            str(address.get("addressCountry") or ""),
        ]).strip()
        if loc:
            dst.setdefault("location", loc)
    rooms = obj.get("numberOfRooms") or obj.get("numberOfBedrooms")
    if rooms:
        dst.setdefault("bedrooms", str(rooms))

# ---------------------------
# Parser principal
# ---------------------------
def parse_listing(html: str, url: str, max_images: int = 3) -> Dict[str, Any]:
    soup = BeautifulSoup(html, "lxml")
    meta = parse_json_ld(soup)
    title = meta.get("title") or text_or_none(soup.find("h1"))
    price = meta.get("price")
    price_str = meta.get("price_str")
    if not price_str:
        price_el = soup.find(["strong", "span"], attrs={"aria-label": re.compile("Preço", re.I)})
        price_str = text_or_none(price_el) or (str(price) if price else "")
    location = meta.get("location") or ""
    if not location:
        breadcrumb = soup.select_one("nav[aria-label='breadcrumb']")
        location = text_or_none(breadcrumb) or ""
    typology   = find_label_value(soup, ["Tipologia"])
    bedrooms   = meta.get("bedrooms") or find_label_value(soup, ["Quartos", "Nº de quartos", "Número de quartos"])
    bathrooms  = find_label_value(soup, ["Casas de banho", "WCs"])
    area       = find_label_value(soup, ["Área bruta", "Área útil", "Área", "Área (m²)", "Área bruta (m²)"])
    description = meta.get("description")
    if not description:
        paragraphs = [p.get_text(" ", strip=True) for p in soup.find_all("p")]
        paragraphs.sort(key=len, reverse=True)
        description = paragraphs[0] if paragraphs else ""
    images: List[str] = []
    for img in soup.find_all("img"):
        src = img.get("src") or img.get("data-src") or img.get("data-lazy")
        if not src:
            continue
        if src.startswith("//"):
            src = "https:" + src
        if src.startswith("http"):
            images.append(src)
        if len(images) >= max_images:
            break
    rec: Dict[str, Any] = {
        "url": url,
        "title": (title or "").strip(),
        "price": (price_str or "").strip(),
        "location": (location or "").strip(),
        "area": (area or "").strip(),
        "typology": (typology or "").strip(),
        "bedrooms": (str(bedrooms) if bedrooms else "").strip(),
        "bathrooms": (str(bathrooms) if bathrooms else "").strip(),
        "description": (description or "").strip(),
    }
    for i, src in enumerate(images, start=1):
        rec[f"image{i}"] = src
    return rec

# ---------------------------
# PowerPoint helpers
# ---------------------------
def add_titlebox(slide, left, top, width, height, text, font_size=28, bold=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame; p.clear()
    run = p.paragraphs[0].add_run(); run.text = text
    run.font.size = Pt(font_size); run.font.bold = bold
    p.paragraphs[0].alignment = PP_ALIGN.LEFT
    return tb

def add_textbox(slide, left, top, width, height, text, font_size=14, color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.clear()
    run = tf.paragraphs[0].add_run(); run.text = text; run.font.size = Pt(font_size)
    if color is not None:
        run.font.color.rgb = color
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    return tb

def build_pptx(df: pd.DataFrame, out_path: Path, brand: str, max_images: int = 3):
    # Load template if present
    if Path(TEMPLATE_FILE).exists():
        prs = Presentation(TEMPLATE_FILE)
    else:
        prs = Presentation()

    # Colors aligned with template_hugosilva
    primary_red = RGBColor(0xdc, 0x35, 0x45)
    secondary_blue = RGBColor(0x0d, 0x6e, 0xfd)

    # Cover (blue background + brand)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    rect = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = secondary_blue; rect.line.fill.background()
    add_titlebox(cover, Inches(0.8), Inches(1.0), Inches(9), Inches(1.2), "Portfólio de Imóveis", 36, True)
    add_textbox(cover, Inches(0.8), Inches(2.2), Inches(9), Inches(0.6), brand, 16, color=RGBColor(255,255,255))

    # Slides por imóvel
    for _, r in df.iterrows():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        # top red bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.35))
        bar.fill.solid(); bar.fill.fore_color.rgb = primary_red; bar.line.fill.background()

        title_line = f"{(r.get('title') or '').strip()} — {(r.get('location') or '').strip()}"
        price_line = (r.get('price') or '').strip()
        add_titlebox(s, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.7), title_line, 24, True)
        add_textbox(s, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.4), price_line, 18, color=primary_red)

        specs = []
        if r.get('typology'): specs.append(f"Tipologia: {r['typology']}")
        if r.get('area'): specs.append(f"Área: {r['area']}")
        if r.get('bedrooms'): specs.append(f"Quartos: {r['bedrooms']}")
        if r.get('bathrooms'): specs.append(f"WCs: {r['bathrooms']}")
        specs.append(f"Link: {r.get('url','')}")
        add_textbox(s, Inches(0.6), Inches(1.8), Inches(5.5), Inches(1.0), "  •  ".join(specs), 12)

        descr = (r.get('description') or '')[:900]
        add_textbox(s, Inches(0.6), Inches(2.7), Inches(5.5), Inches(3.5), descr, 12)

        col = 0
        for i in range(1, max_images+1):
            key = f"image{i}"
            if not r.get(key): 
                continue
            try:
                img_bytes = requests.get(r[key], timeout=12).content
                s.shapes.add_picture(io.BytesIO(img_bytes), Inches(6.3), Inches(1.8 + col*2.0), width=Inches(3.0))
                col += 1
            except Exception:
                pass

    prs.save(out_path)

# ---------------------------
# Main
# ---------------------------
async def run(args):
    df_urls = pd.read_csv(args.input)
    rows = []
    for _, row in df_urls.iterrows():
        url = str(row["url"]).strip()
        if not url or not url.startswith("http"):
            continue
        try:
            html = await fetch_html(url, render=args.render)
            rec = parse_listing(html, url, max_images=args.max_images)
            rows.append(rec)
            print(f"[OK] {url} → {rec.get('title','(sem título)')[:80]}")
        except Exception as e:
            print(f"[ERRO] {url}: {e}")
        time.sleep(args.delay)

    if not rows:
        print("Nenhum registo extraído.")
        return

    max_imgs = max([len([k for k in r.keys() if k.startswith("image")]) for r in rows])
    for r in rows:
        for i in range(1, max_imgs+1):
            r.setdefault(f"image{i}", "")

    out_csv = Path(args.output)
    pd.DataFrame(rows).to_csv(out_csv, index=False, quoting=csv.QUOTE_NONNUMERIC)

    out_pptx = Path(args.pptx)
    build_pptx(pd.DataFrame(rows), out_pptx, brand=args.brand, max_images=args.max_images)

    print(f"Feito! CSV: {out_csv.resolve()} | PPTX: {out_pptx.resolve()}")

def parse_args():
    p = argparse.ArgumentParser(description="Extrai dados do Imovirtual → CSV + PPTX (template Hugo Silva)")
    p.add_argument("--input", required=True, help="CSV com coluna 'url'")
    p.add_argument("--output", default="dados.csv", help="CSV de saída com dados")
    p.add_argument("--pptx", default="Apresentacoes_Imoveis.pptx", help="Nome do PowerPoint de saída")
    p.add_argument("--brand", default="Hugo Silva | Consultor Imobiliário · RE/MAX Oceanus — Choose your dream. Live in it.",
                   help="Assinatura/copy para a apresentação")
    p.add_argument("--delay", type=float, default=2.0, help="Pausa (segundos) entre URLs")
    p.add_argument("--render", choices=["load", "domcontentloaded", "networkidle"],
                   default="networkidle", help="Playwright wait_until")
    p.add_argument("--max-images", type=int, default=3, help="Máximo de imagens por imóvel")
    return p.parse_args()

if __name__ == "__main__":
    args = parse_args()
    asyncio.run(run(args))
